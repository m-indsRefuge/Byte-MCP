"""Safe, bounded text extraction for V1 file types."""
from __future__ import annotations

import csv
import io
import zipfile
from pathlib import Path

from .errors import UnsupportedFileError

TEXT_EXTENSIONS = frozenset(
    {
        ".txt",
        ".md",
        ".json",
        ".jsonl",
        ".csv",
        ".tsv",
        ".log",
        ".xml",
        ".toml",
        ".ini",
        ".cfg",
        ".yaml",
        ".yml",
        ".py",
        ".ps1",
        ".js",
        ".jsx",
        ".ts",
        ".tsx",
        ".html",
        ".css",
        ".sql",
        ".bat",
        ".cmd",
        ".sh",
        ".java",
        ".c",
        ".h",
        ".cpp",
        ".hpp",
    }
)
DOCUMENT_EXTENSIONS = frozenset(
    {".pdf", ".docx", ".xlsx", ".pptx", ".zip"}
)
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | DOCUMENT_EXTENSIONS


def _bounded(text: str, max_chars: int) -> tuple[str, bool]:
    return text[:max_chars], len(text) > max_chars


def _read_text(path: Path) -> str:
    raw = path.read_bytes()
    decoded: str | None = None

    for encoding in ("utf-8-sig", "utf-16", "cp1252"):
        try:
            decoded = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue

    if decoded is None:
        decoded = raw.decode("utf-8", errors="replace")

    return decoded.replace("\r\n", "\n").replace("\r", "\n")


def extract_file(
    path: Path,
    max_chars: int,
) -> tuple[str, bool, str]:
    suffix = path.suffix.casefold()

    if suffix in TEXT_EXTENSIONS:
        text, truncated = _bounded(_read_text(path), max_chars)
        return text, truncated, "text"

    if suffix == ".pdf":
        from pypdf import PdfReader

        text = "\n\n".join(
            page.extract_text() or ""
            for page in PdfReader(str(path)).pages
        )
        text, truncated = _bounded(text, max_chars)
        return text, truncated, "pdf"

    if suffix == ".docx":
        from docx import Document

        document = Document(str(path))
        lines = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            lines.extend(
                " | ".join(str(cell.text) for cell in row.cells)
                for row in table.rows
            )
        text, truncated = _bounded("\n".join(lines), max_chars)
        return text, truncated, "docx"

    if suffix == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(
            path,
            read_only=True,
            data_only=True,
        )
        output = io.StringIO()
        writer = csv.writer(output)
        truncated = False

        for sheet in workbook.worksheets:
            output.write(f"\n## Sheet: {sheet.title}\n")
            for row in sheet.iter_rows(values_only=True):
                writer.writerow(
                    ["" if value is None else value for value in row]
                )
                if output.tell() >= max_chars:
                    truncated = True
                    break
            if truncated:
                break

        workbook.close()
        return output.getvalue()[:max_chars], truncated, "xlsx"

    if suffix == ".pptx":
        from pptx import Presentation

        lines: list[str] = []
        for index, slide in enumerate(
            Presentation(str(path)).slides,
            start=1,
        ):
            lines.append(f"\n## Slide {index}")
            lines.extend(
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )

        text, truncated = _bounded("\n".join(lines), max_chars)
        return text, truncated, "pptx"

    if suffix == ".zip":
        lines: list[str] = []
        with zipfile.ZipFile(path) as archive:
            for index, item in enumerate(
                archive.infolist(),
                start=1,
            ):
                if index > 500:
                    lines.append(
                        "... archive listing truncated at 500 entries ..."
                    )
                    break
                lines.append(
                    f"{item.filename}\t{item.file_size} bytes"
                )

        text, truncated = _bounded("\n".join(lines), max_chars)
        return text, truncated, "zip-listing"

    raise UnsupportedFileError(
        f"No V1 extractor is registered for "
        f"{suffix or '<no extension>'}."
    )
